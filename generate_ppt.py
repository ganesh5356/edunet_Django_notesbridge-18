import os
import sys
import subprocess

def install_pptx():
    try:
        import pptx
    except ImportError:
        print("Installing python-pptx library to create the PowerPoint...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])

install_pptx()

from pptx import Presentation

def create_presentation():
    # Initialize Presentation
    prs = Presentation()

    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "NotesBridge"
    subtitle.text = "Connecting Seniors' Knowledge to Juniors' Success!\n\nCollege: Kishkinda University"

    # 2. Disclaimer
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    title_shape.text = "Disclaimer"
    body_shape.text_frame.text = "The content is curated from online/offline resources and used for educational purpose only."

    # 3. Problem Statement & Introduction
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    title_shape.text = "The Problem & Introduction"
    tf = body_shape.text_frame
    tf.text = "The Problem Statement:"
    
    p = tf.add_paragraph()
    p.text = "• Fragmented data and scattered study resources"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Lack of organized mentorship from seniors to juniors"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Difficulty in resolving academic doubts reliably"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Introduction:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• A centralized web-based platform for academic collaboration"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Intelligent role-based dashboards (Juniors vs Seniors)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Built-in real-time analytics for engagement tracking"
    p.level = 1

    # 4. Purpose & Solution
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "NotesBridge: Purpose & Solution"
    tf = slide.placeholders[1].text_frame
    tf.text = "Key Features:"
    
    features = [
        "Role-Specific Access & Tailored Dashboards",
        "Secure Document Upload & Robust Local Preview",
        "Community Doubts Forum & Resolution System",
        "Advanced Analytics Dashboard (Chart.js)",
        "Resource Bookmarking & Download Tracking",
        "Modern Responsive UI across all devices"
    ]
    for feat in features:
        p = tf.add_paragraph()
        p.text = "• " + feat
        p.level = 1

    # 5. Architecture
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "System Architecture"
    tf = slide.placeholders[1].text_frame
    tf.text = "Frontend (User Interface)"
    p = tf.add_paragraph()
    p.text = "• HTML5, CSS3 (Custom Design System), JavaScript, Chart.js"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Backend (Django Framework)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• REST API, User Auth, Resource Management Logic"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Database & Services"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• Relational Database (SQLite/PostgreSQL)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Media Storage for Notes, PDFs, and Profile Images"
    p.level = 1

    # 6. Technical Roadmap
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Technical Roadmap"
    tf = slide.placeholders[1].text_frame
    steps = [
        "1. Core Assembly: Backend setup, Profiles, Auth system",
        "2. Resource Engine: File uploads, Previews, Base Storage",
        "3. Community: Doubts forum, Replies, Bookmarking logic",
        "4. UI Refinement: Glassmorphism, Responsive CSS structure",
        "5. Dashboard Live: Role-based logic, Visual analytics",
        "6. Final Polish: Security, Mobile fixes, Footer customization"
    ]
    tf.text = steps[0]
    for step in steps[1:]:
        p = tf.add_paragraph()
        p.text = step
        p.level = 0

    # 7. Platform Screenshots
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Platform Screenshots"
    tf = slide.placeholders[1].text_frame
    tf.text = "(You can drag and drop your project screenshots onto this slide)"
    p = tf.add_paragraph()
    p.text = "• Recommend adding the Landing Page with Dynamic Hero Section"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Recommend adding the Advanced Statistics Dashboard"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Recommend adding the Document Preview & Detail views"
    p.level = 1

    # 8. Conclusion
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Conclusion"
    tf = slide.placeholders[1].text_frame
    conc = [
        "1. Secure & Centralized: Unified platform tailored for university students.",
        "2. Insightful Intelligence: Automated tracking of downloads and engagement.",
        "3. Community Driven: Bridges the gap across batches, fostering collaborative peer learning via shared notes and Q&A."
    ]
    tf.text = conc[0]
    for c in conc[1:]:
        p = tf.add_paragraph()
        p.text = c
        p.level = 0

    # 9. References
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "References"
    tf = slide.placeholders[1].text_frame
    refs = [
        "Django Documentation (djangoproject.com)",
        "MDN Web Docs",
        "Chart.js Documentation",
        "Modern Responsive Web Design Practices"
    ]
    tf.text = refs[0]
    for r in refs[1:]:
        p = tf.add_paragraph()
        p.text = r
        p.level = 0

    # Save the presentation
    output_path = os.path.abspath("NotesBridge_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved successfully at: {output_path}")

if __name__ == '__main__':
    create_presentation()
